import warnings
warnings.filterwarnings("ignore")

import os
import random
from collections import deque
from PIL import Image
import numpy as np
from pptx import Presentation
from pptx.util import Emu

BASE_DIR = r"C:\Users\北興\Desktop\claude-ai"
STICKER_PATH = os.path.join(BASE_DIR, "sticker.png")
PPTX_PATH    = os.path.join(BASE_DIR, "木蘭詩簡報.pptx")
OUTPUT_PATH  = os.path.join(BASE_DIR, "木蘭詩簡報_with_stickers.pptx")
CHAR_DIR     = os.path.join(BASE_DIR, "sticker_chars")
os.makedirs(CHAR_DIR, exist_ok=True)


# ── 1. 偵測角色邊界（grid-based BFS 連通區域） ────────────────────────────
def find_characters(img_path, cell_size=15, density=0.08, min_cells=10, padding=20):
    img = Image.open(img_path).convert("RGBA")
    arr = np.array(img)
    h, w = arr.shape[:2]

    # 非白像素遮罩
    r, g, b = arr[:,:,0], arr[:,:,1], arr[:,:,2]
    non_white = ~((r > 240) & (g > 240) & (b > 240))

    # 格子密度門檻（需達 density 比例才算有內容）
    grid_h = h // cell_size
    grid_w = w // cell_size
    grid = np.zeros((grid_h, grid_w), dtype=bool)
    for gy in range(grid_h):
        for gx in range(grid_w):
            chunk = non_white[gy*cell_size:(gy+1)*cell_size,
                              gx*cell_size:(gx+1)*cell_size]
            if chunk.mean() >= density:
                grid[gy, gx] = True

    # BFS 找連通區域
    visited = np.zeros_like(grid, dtype=bool)
    components = []
    for sy in range(grid_h):
        for sx in range(grid_w):
            if grid[sy, sx] and not visited[sy, sx]:
                q = deque([(sy, sx)])
                visited[sy, sx] = True
                cells = [(sy, sx)]
                while q:
                    cy, cx = q.popleft()
                    for dy in (-1, 0, 1):
                        for dx in (-1, 0, 1):
                            ny, nx = cy+dy, cx+dx
                            if 0 <= ny < grid_h and 0 <= nx < grid_w:
                                if grid[ny, nx] and not visited[ny, nx]:
                                    visited[ny, nx] = True
                                    q.append((ny, nx))
                                    cells.append((ny, nx))
                if len(cells) >= min_cells:
                    components.append(cells)

    # 轉回像素座標 + padding
    bboxes = []
    for cells in components:
        ys = [c[0] for c in cells]
        xs = [c[1] for c in cells]
        x1 = max(0, min(xs)*cell_size - padding)
        y1 = max(0, min(ys)*cell_size - padding)
        x2 = min(w, (max(xs)+1)*cell_size + padding)
        y2 = min(h, (max(ys)+1)*cell_size + padding)
        bboxes.append((x1, y1, x2, y2))

    bboxes.sort(key=lambda b: (b[1], b[0]))
    return img, bboxes


# ── 2. 裁切並讓白底透明 ─────────────────────────────────────────────────
def crop_and_transparent(img, bbox, out_path, white_thresh=240):
    x1, y1, x2, y2 = bbox
    crop = img.crop((x1, y1, x2, y2)).copy()
    arr  = np.array(crop)
    r, g, b = arr[:,:,0], arr[:,:,1], arr[:,:,2]
    white_mask = (r > white_thresh) & (g > white_thresh) & (b > white_thresh)
    arr[:,:,3] = np.where(white_mask, 0, 255)
    Image.fromarray(arr).save(out_path, "PNG")


# ── 3. 主流程 ────────────────────────────────────────────────────────────
print("=== 裁切角色 ===")
sheet_img, bboxes = find_characters(STICKER_PATH)
print(f"偵測到 {len(bboxes)} 個角色")

char_paths = []
for i, bbox in enumerate(bboxes):
    path = os.path.join(CHAR_DIR, f"char_{i+1:02d}.png")
    crop_and_transparent(sheet_img, bbox, path)
    w = bbox[2]-bbox[0]; h = bbox[3]-bbox[1]
    print(f"  ✅ char_{i+1:02d}.png  ({w}×{h}px)")
    char_paths.append(path)

print(f"\n共裁切 {len(char_paths)} 個角色，存到 sticker_chars/\n")

# ── 4. 加入簡報 ──────────────────────────────────────────────────────────
print("=== 加入簡報 ===")
prs      = Presentation(PPTX_PATH)
slide_w  = prs.slide_width
slide_h  = prs.slide_height
margin   = int(slide_w * 0.02)
sticker_w = int(slide_w * 0.13)   # 約 13% 簡報寬度

random.seed(42)
used = []

for i, slide in enumerate(prs.slides):
    # 每頁用不同角色（循環不重複直到用完）
    if not used or len(used) == len(char_paths):
        pool = char_paths.copy()
        random.shuffle(pool)
        used = pool.copy()
    char_path = used.pop(0)

    # 依比例算高度
    cimg = Image.open(char_path)
    cw, ch = cimg.size
    sticker_h = int(sticker_w * ch / cw)

    left = slide_w - sticker_w - margin
    top  = slide_h - sticker_h - margin

    slide.shapes.add_picture(char_path, left, top, sticker_w, sticker_h)
    print(f"  📄 第 {i+1:>2d} 頁 → {os.path.basename(char_path)}")

prs.save(OUTPUT_PATH)
print(f"\n✅ 完成！→ {OUTPUT_PATH}")
