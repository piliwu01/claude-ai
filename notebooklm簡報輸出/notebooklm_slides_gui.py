import warnings
warnings.filterwarnings("ignore")

import os, sys, json, random, threading, subprocess, time
from collections import deque
import tkinter as tk
from tkinter import ttk, filedialog
from tkinter.scrolledtext import ScrolledText

import numpy as np
from PIL import Image
from pptx import Presentation
from pptx.util import Emu

# ── 16 種簡報風格 ─────────────────────────────────────────────────────────
STYLES = [
    ("現代報紙",    "Modern newspaper style, Swiss/Bauhaus design, ultra-massive headlines occupying 30-50% of slide, extreme 10:1 size ratio between headline and body text, white or cool gray background, jet black text, electric yellow or alert red accents"),
    ("極簡建築",    "Sharp-edged minimalism, grid-based layout, top-left navigation, generous negative space, monochrome design, light gray or white base, jet black text and accents"),
    ("黃黑時尚",    "Yellow and black editorial style, fashion magazine aesthetic, dynamic serif fonts, yellow background, black text, handwritten or sticker elements"),
    ("橙白創意",    "Black and orange creative agency style, white background, blood orange accents, dynamic photography composition"),
    ("漫畫教育",    "Manga comic style, story panels layout, educational narrative storytelling, visual sequential art format for making complex information relatable"),
    ("粉色街頭",    "Pink street-style, pop and deformed illustrations with thick outlines, soft squishy image shapes, pink background, white and black text, flat colors"),
    ("科技像素",    "Digital neo pop pixel-infographic editorial style, grid paper background, 8-16 bit pixel art icons, modular colored blocks, hot pink, bright yellow, cyan/light blue, black and cream palette, tech and developer aesthetic"),
    ("水彩藝術",    "Royal blue and red watercolor style, wet watercolor technique, fine art gallery aesthetic, artistic and classical mood"),
    ("霓虹科技",    "Tech art neon style, neon color palette, contemporary technology and art fusion aesthetic"),
    ("雕塑波普",    "Sculpture vaporwave style, classical sculpture imagery mixed with vaporwave aesthetics, art history meets contemporary culture"),
    ("Premium產品", "Studio mockup premium style, 3D product mockups, jet black UI screens, white or light gray background, electric purple (#8D59E9) and acid yellow (#EBE021) accents, Apple product aesthetic"),
    ("運動活力",    "Sports athletic energy style, high-energy dynamic layout, bold typography, motivational and powerful visual composition"),
    ("研討會極簡",  "White background, black text, red accent color, sans-serif font, dynamic typography, high-sensibility minimal design for seminar use"),
    ("雜誌可愛",    "Mature-cute editorial magazine style, large cutout photo centered, asymmetrical speech bubbles in margins, numbered sections, matte dusty pink background, hand-drawn speech bubbles, sophisticated headings with casual handwritten fonts"),
    ("明朝手寫",    "Yellow background, black text, large dynamic modern serif font placement, pop and chic touches like handwriting or stickers scattered throughout, bold fashion magazine layout"),
    ("扁平插畫",    "Flat color illustration with slightly distorted character proportions and thick outlines, limited palette of gentle tones mixed with white (maximum 3 colors), solid single-color background, soft muted coloring with bold line work"),
]

# ── 色彩 ──────────────────────────────────────────────────────────────────
BG      = "#f0f4f8"
CARD    = "#ffffff"
BTN_B   = "#2563eb"
BTN_G   = "#16a34a"
T_MAIN  = "#1e293b"
T_SUB   = "#64748b"
LOG_BG  = "#1e293b"
LOG_FG  = "#94a3b8"
C_OK    = "#4ade80"
C_ERR   = "#f87171"
C_GRAY  = "#475569"
FONT    = "微軟正黑體"


class App(tk.Tk):
    def __init__(self):
        super().__init__()
        self.title("NotebookLM 簡報製作器 v1.0")
        self.configure(bg=BG)
        self.resizable(False, False)
        self._output_dir = None
        self._running = False
        self._build_ui()
        self.after(100, self._center)

    def _center(self):
        self.update_idletasks()
        w, h = self.winfo_width(), self.winfo_height()
        sw, sh = self.winfo_screenwidth(), self.winfo_screenheight()
        self.geometry(f"+{(sw-w)//2}+{(sh-h)//2}")

    # ── UI ────────────────────────────────────────────────────────────────
    def _build_ui(self):
        PAD = dict(padx=14, pady=6)

        # 標題
        tk.Label(self, text="📊  NotebookLM 簡報製作器",
                 font=(FONT, 16, "bold"), bg=BG, fg=T_MAIN).pack(padx=14, pady=(16, 2))
        tk.Label(self, text="PDF → NotebookLM 生成簡報 → 加入個人貼圖",
                 font=(FONT, 10), bg=BG, fg=T_SUB).pack(pady=(0, 8))

        ttk.Separator(self).pack(fill="x", padx=14)

        # ① PDF
        frm_pdf = tk.LabelFrame(self, text="  ① 選擇 PDF  ",
                                font=(FONT, 10, "bold"), bg=CARD, fg="#334155",
                                relief="groove", bd=1, padx=10, pady=8)
        frm_pdf.pack(fill="x", padx=14, pady=(10, 4))
        self.var_pdf = tk.StringVar()
        tk.Entry(frm_pdf, textvariable=self.var_pdf, width=52,
                 font=(FONT, 10), relief="solid", bd=1,
                 state="readonly", readonlybackground="#f8fafc").pack(
                     side="left", ipady=4, padx=(0, 8))
        tk.Button(frm_pdf, text="選擇…", command=self._pick_pdf,
                  bg=BTN_B, fg="white", font=(FONT, 10, "bold"),
                  relief="flat", padx=10, cursor="hand2").pack(side="left")

        # ② 輸出資料夾
        frm_out = tk.LabelFrame(self, text="  ② 輸出資料夾  ",
                                font=(FONT, 10, "bold"), bg=CARD, fg="#334155",
                                relief="groove", bd=1, padx=10, pady=8)
        frm_out.pack(fill="x", padx=14, pady=4)
        self.var_outdir = tk.StringVar()
        tk.Entry(frm_out, textvariable=self.var_outdir, width=52,
                 font=(FONT, 10), relief="solid", bd=1,
                 state="readonly", readonlybackground="#f8fafc").pack(
                     side="left", ipady=4, padx=(0, 8))
        tk.Button(frm_out, text="選擇…", command=self._pick_outdir,
                  bg=BTN_B, fg="white", font=(FONT, 10, "bold"),
                  relief="flat", padx=10, cursor="hand2").pack(side="left")
        tk.Label(frm_out, text="（未選則與 PDF 同資料夾）",
                 font=(FONT, 9), bg=CARD, fg=T_SUB).pack(side="left", padx=(8, 0))

        # ③ 貼圖（可略）
        frm_stk = tk.LabelFrame(self, text="  ③ 選擇貼圖（可略過）  ",
                                font=(FONT, 10, "bold"), bg=CARD, fg="#334155",
                                relief="groove", bd=1, padx=10, pady=8)
        frm_stk.pack(fill="x", padx=14, pady=4)
        self.var_sticker = tk.StringVar()
        tk.Entry(frm_stk, textvariable=self.var_sticker, width=52,
                 font=(FONT, 10), relief="solid", bd=1,
                 state="readonly", readonlybackground="#f8fafc").pack(
                     side="left", ipady=4, padx=(0, 8))
        tk.Button(frm_stk, text="選擇…", command=self._pick_sticker,
                  bg=BTN_B, fg="white", font=(FONT, 10, "bold"),
                  relief="flat", padx=10, cursor="hand2").pack(side="left")
        tk.Button(frm_stk, text="清除", command=lambda: self.var_sticker.set(""),
                  bg="#e2e8f0", font=(FONT, 9),
                  relief="flat", padx=8, cursor="hand2").pack(side="left", padx=(6, 0))

        # ④ 風格選擇
        frm_style = tk.LabelFrame(self, text="  ④ 選擇簡報風格  ",
                                  font=(FONT, 10, "bold"), bg=CARD, fg="#334155",
                                  relief="groove", bd=1, padx=10, pady=8)
        frm_style.pack(fill="x", padx=14, pady=4)

        self.var_style = tk.IntVar(value=4)   # 預設漫畫教育（index 4）
        for i, (name, _) in enumerate(STYLES):
            row, col = divmod(i, 2)
            rb = tk.Radiobutton(frm_style, text=f"{i+1:02d}. {name}",
                                variable=self.var_style, value=i,
                                font=(FONT, 10), bg=CARD, fg=T_MAIN,
                                activebackground=CARD, selectcolor="#dbeafe",
                                cursor="hand2")
            rb.grid(row=row, column=col, sticky="w", padx=8, pady=2)

        # 執行按鈕
        self.btn_run = tk.Button(self, text="▶  開始製作",
                                 command=self._run,
                                 bg=BTN_G, fg="white",
                                 font=(FONT, 13, "bold"),
                                 relief="flat", padx=30, pady=10,
                                 cursor="hand2")
        self.btn_run.pack(pady=(12, 4))

        # 進度條
        self.progress = ttk.Progressbar(self, mode="indeterminate", length=460)
        self.progress.pack(padx=14, pady=(0, 6))

        ttk.Separator(self).pack(fill="x", padx=14)

        # 執行記錄
        self.log = ScrolledText(self, height=12, width=62,
                                font=("Consolas", 10), bg=LOG_BG,
                                fg=LOG_FG, relief="flat",
                                state="disabled", padx=6, pady=6)
        self.log.tag_config("ok",   foreground=C_OK)
        self.log.tag_config("err",  foreground=C_ERR)
        self.log.tag_config("sep",  foreground=C_GRAY)
        self.log.tag_config("info", foreground=LOG_FG)
        self.log.pack(padx=14, pady=(6, 4))

        # 開啟資料夾按鈕（完成後顯示）
        self.btn_open = tk.Button(self, text="📂  開啟輸出資料夾",
                                  command=self._open_folder,
                                  bg="#0f172a", fg=C_OK,
                                  font=(FONT, 10, "bold"),
                                  relief="flat", padx=16, pady=6,
                                  cursor="hand2")
        # 先不 pack，完成後才顯示

        tk.Label(self, text="⚠  生成時間約 5～45 分鐘，請耐心等候",
                 font=(FONT, 9), bg=BG, fg=T_SUB).pack(pady=(0, 14))

    # ── 事件 ──────────────────────────────────────────────────────────────
    def _pick_pdf(self):
        path = filedialog.askopenfilename(
            title="選擇 PDF",
            filetypes=[("PDF 文件", "*.pdf"), ("所有檔案", "*.*")])
        if path:
            self.var_pdf.set(path)

    def _pick_outdir(self):
        path = filedialog.askdirectory(title="選擇輸出資料夾")
        if path:
            self.var_outdir.set(path)

    def _pick_sticker(self):
        path = filedialog.askopenfilename(
            title="選擇貼圖",
            filetypes=[("圖片", "*.png *.jpg *.jpeg"), ("所有檔案", "*.*")])
        if path:
            self.var_sticker.set(path)

    def _run(self):
        if self._running:
            return
        if not self.var_pdf.get():
            self._log("❌  請先選擇 PDF 檔案", "err")
            return
        self._running = True
        self.btn_run.config(state="disabled", text="製作中…")
        self.progress.start(12)
        self.btn_open.pack_forget()
        threading.Thread(target=self._process, daemon=True).start()

    def _open_folder(self):
        if self._output_dir and os.path.isdir(self._output_dir):
            os.startfile(self._output_dir)

    # ── 背景執行緒 ────────────────────────────────────────────────────────
    def _process(self):
        pdf_path     = self.var_pdf.get()
        sticker_path = self.var_sticker.get()
        style_idx    = self.var_style.get()
        style_name, style_prompt = STYLES[style_idx]

        out_dir  = self.var_outdir.get() or os.path.dirname(pdf_path)
        pdf_name = os.path.splitext(os.path.basename(pdf_path))[0]
        char_dir = os.path.join(out_dir, "sticker_chars")

        self._log("=" * 50, "sep")
        self._log(f"📄  PDF：{os.path.basename(pdf_path)}", "info")
        self._log(f"🎨  風格：{style_name}", "info")
        self._log("=" * 50, "sep")

        # 1. 建立筆記本
        self._log("📋  建立 NotebookLM 筆記本…")
        ok, out = self._cmd(["notebooklm", "create", pdf_name, "--json"])
        if not ok:
            return self._fail(f"建立筆記本失敗：{out}")
        try:
            notebook_id = json.loads(out)["notebook"]["id"]
        except Exception:
            return self._fail(f"解析筆記本 ID 失敗：{out}")
        self._log(f"✅  筆記本建立：{notebook_id[:8]}…", "ok")

        self._cmd(["notebooklm", "use", notebook_id])

        # 2. 上傳 PDF
        self._log("📤  上傳 PDF…")
        ok, out = self._cmd(["notebooklm", "source", "add", pdf_path, "--json"])
        if not ok:
            return self._fail(f"上傳 PDF 失敗：{out}")
        try:
            source_id = json.loads(out)["source"]["id"]
        except Exception:
            return self._fail(f"解析 source ID 失敗：{out}")
        self._log(f"✅  上傳成功：{source_id[:8]}…", "ok")

        # 3. 等待 PDF 處理
        self._log("⏳  等待 PDF 處理（最多 5 分鐘）…")
        ok, out = self._cmd(["notebooklm", "source", "wait",
                             source_id, "--timeout", "300"])
        if not ok:
            return self._fail(f"PDF 處理失敗：{out}")
        self._log("✅  PDF 處理完成", "ok")

        # 4. 生成簡報
        self._log(f"🎨  開始生成簡報（{style_name}）…")
        full_prompt = f"Create exactly 30 slides. {style_prompt}"
        ok, out = self._cmd(["notebooklm", "generate", "slide-deck",
                             full_prompt, "--language", "zh_Hant", "--json"])
        if not ok:
            return self._fail(f"生成指令失敗：{out}")
        try:
            artifact_id = json.loads(out)["task_id"]
        except Exception:
            return self._fail(f"解析 artifact ID 失敗：{out}")
        self._log(f"✅  生成任務已啟動：{artifact_id[:8]}…", "ok")

        # 5. 輪詢等待完成
        self._log("⏳  等待生成完成（每 60 秒檢查一次）…")
        final_status = self._poll_artifact(notebook_id, artifact_id)
        if final_status != "completed":
            return self._fail("簡報生成失敗或逾時，請稍後在 NotebookLM 網頁確認")

        # 6. 下載 PPTX
        pptx_base  = os.path.join(out_dir, f"{pdf_name}_簡報.pptx")
        final_pptx = os.path.join(out_dir, f"{pdf_name}_簡報_貼圖版.pptx")
        self._log("📥  下載 PPTX…")
        ok, out = self._cmd(["notebooklm", "download", "slide-deck",
                             pptx_base, "--format", "pptx",
                             "-a", artifact_id, "-n", notebook_id])
        if not ok:
            return self._fail(f"下載失敗：{out}")
        self._log(f"✅  PPTX 下載完成", "ok")

        # 7. 貼圖（可略）
        if sticker_path and os.path.isfile(sticker_path):
            self._log("✂️   裁切角色貼圖…")
            try:
                os.makedirs(char_dir, exist_ok=True)
                char_paths = self._extract_stickers(sticker_path, char_dir)
                self._log(f"✅  裁切 {len(char_paths)} 個角色", "ok")

                self._log("🖼️   加入貼圖到每頁…")
                self._add_stickers(pptx_base, char_paths, final_pptx)
                self._log(f"✅  貼圖版：{os.path.basename(final_pptx)}", "ok")
                output_file = final_pptx
            except Exception as e:
                self._log(f"⚠️   貼圖步驟失敗（{e}），改用無貼圖版", "err")
                output_file = pptx_base
        else:
            output_file = pptx_base

        self._output_dir = out_dir
        self._log("=" * 50, "sep")
        self._log(f"🎉  完成！→ {os.path.basename(output_file)}", "ok")
        self._log("=" * 50, "sep")
        self._done()

    # ── 輪詢 artifact 狀態 ────────────────────────────────────────────────
    def _poll_artifact(self, notebook_id, artifact_id, max_minutes=60):
        for minute in range(1, max_minutes + 1):
            time.sleep(60)
            ok, out = self._cmd(["notebooklm", "artifact", "list",
                                 "--notebook", notebook_id, "--json"])
            if ok:
                try:
                    data = json.loads(out)
                    for a in data.get("artifacts", []):
                        if a["id"] == artifact_id:
                            status = a.get("status", "unknown")
                            if status == "completed":
                                return "completed"
                            if status not in ("in_progress", "pending"):
                                return status
                except Exception:
                    pass
            self._log(f"  仍在生成中… ({minute} 分鐘)", "info")
        return "timeout"

    # ── 裁切貼圖 ──────────────────────────────────────────────────────────
    def _extract_stickers(self, img_path, out_dir,
                          cell_size=15, density=0.08, min_cells=10, padding=20):
        img = Image.open(img_path).convert("RGBA")
        arr = np.array(img)
        h, w = arr.shape[:2]

        r, g, b = arr[:,:,0], arr[:,:,1], arr[:,:,2]
        non_white = ~((r > 240) & (g > 240) & (b > 240))

        grid_h, grid_w = h // cell_size, w // cell_size
        grid = np.zeros((grid_h, grid_w), dtype=bool)
        for gy in range(grid_h):
            for gx in range(grid_w):
                chunk = non_white[gy*cell_size:(gy+1)*cell_size,
                                  gx*cell_size:(gx+1)*cell_size]
                if chunk.mean() >= density:
                    grid[gy, gx] = True

        visited = np.zeros_like(grid, dtype=bool)
        components = []
        for sy in range(grid_h):
            for sx in range(grid_w):
                if grid[sy, sx] and not visited[sy, sx]:
                    q = deque([(sy, sx)])
                    visited[sy, sx] = True
                    cells = [(sy, sx)]
                    while q:
                        cy, cx = q.popleft()
                        for dy in (-1, 0, 1):
                            for dx in (-1, 0, 1):
                                ny, nx = cy+dy, cx+dx
                                if (0 <= ny < grid_h and 0 <= nx < grid_w
                                        and grid[ny, nx] and not visited[ny, nx]):
                                    visited[ny, nx] = True
                                    q.append((ny, nx))
                                    cells.append((ny, nx))
                    if len(cells) >= min_cells:
                        components.append(cells)

        components.sort(key=lambda cells: (
            min(c[0] for c in cells), min(c[1] for c in cells)))

        char_paths = []
        for idx, cells in enumerate(components):
            ys = [c[0] for c in cells]
            xs = [c[1] for c in cells]
            x1 = max(0, min(xs)*cell_size - padding)
            y1 = max(0, min(ys)*cell_size - padding)
            x2 = min(w, (max(xs)+1)*cell_size + padding)
            y2 = min(h, (max(ys)+1)*cell_size + padding)

            crop = img.crop((x1, y1, x2, y2)).copy()
            ca   = np.array(crop)
            rr, gg, bb = ca[:,:,0], ca[:,:,1], ca[:,:,2]
            wm = (rr > 240) & (gg > 240) & (bb > 240)
            ca[:,:,3] = np.where(wm, 0, 255)
            path = os.path.join(out_dir, f"char_{idx+1:02d}.png")
            Image.fromarray(ca).save(path, "PNG")
            char_paths.append(path)

        return char_paths

    # ── 加貼圖到 PPTX ─────────────────────────────────────────────────────
    def _add_stickers(self, pptx_path, char_paths, output_path):
        prs    = Presentation(pptx_path)
        sw     = prs.slide_width
        sh     = prs.slide_height
        sw_px  = int(sw * 0.13)
        margin = int(sw * 0.02)

        random.seed(42)
        pool   = char_paths.copy()
        random.shuffle(pool)
        used   = []

        for slide in prs.slides:
            if not used or len(used) == len(char_paths):
                pool = char_paths.copy()
                random.shuffle(pool)
                used = pool.copy()
            cp = used.pop(0)
            ci = Image.open(cp)
            cw, ch = ci.size
            sh_px = int(sw_px * ch / cw)
            left  = sw - sw_px - margin
            top   = sh - sh_px - margin
            slide.shapes.add_picture(cp, left, top, sw_px, sh_px)

        prs.save(output_path)

    # ── 工具 ──────────────────────────────────────────────────────────────
    def _cmd(self, args):
        try:
            r = subprocess.run(args, capture_output=True,
                               text=True, encoding="utf-8", errors="replace")
            out = r.stdout.strip()
            err = r.stderr.strip()
            return (r.returncode == 0), (out or err)
        except Exception as e:
            return False, str(e)

    def _log(self, msg, tag="info"):
        def _append():
            self.log.config(state="normal")
            self.log.insert("end", msg + "\n", tag)
            self.log.see("end")
            self.log.config(state="disabled")
        self.after(0, _append)

    def _fail(self, msg):
        self._log(f"❌  {msg}", "err")
        self.after(0, self._reset)

    def _done(self):
        self.after(0, lambda: (
            self._reset(),
            self.btn_open.pack(pady=(2, 14))
        ))

    def _reset(self):
        self._running = False
        self.progress.stop()
        self.btn_run.config(state="normal", text="▶  開始製作")


if __name__ == "__main__":
    App().mainloop()
